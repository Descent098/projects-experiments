from json import load

from pptx import Presentation
from jinja2 import Environment, FileSystemLoader

def read_powerpoint(file_path):
    # Load the PowerPoint presentation
    prs = Presentation(file_path)
    
    content = []
    
    # Iterate through each slide in the presentation
    for slide in prs.slides:
        print(f"Slide {prs.slides.index(slide) + 1}:")
        
        
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(f"\tShape: {shape.text}")
                content.append(shape.text.replace("’","'"))
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    print(f"\tParagraph: {paragraph.text}")
                    content.append(paragraph.text)
    return content
                    

with open("config.json", "r") as conf:
    config = load(conf)
    
# Path to your PowerPoint file
pptx_file = "example.pptx"

# Call the function to read the PowerPoint file
content = read_powerpoint(pptx_file)

print(content)

# Create a Jinja2 environment with a file system loader
env = Environment(loader=FileSystemLoader("templates"))


# Load the template from a file
homepage = env.get_template("index.jinja")

# Render the template with the provided data
rendered_homepage = homepage.render(config=config, slide_content=content)

with open("index.html", "w+") as out:
    out.write(rendered_homepage)

